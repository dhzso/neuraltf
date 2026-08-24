#!/usr/bin/env python
'''State-of-the-art NeuralTF PowerPoint generator.

* Uses the Deepanshu master template for styling.
* One figure per slide (plus brief explanatory caption).
* Slides: Title, Outline, then each figure from the main set.
* Additional Dirichlet robustness slides are grouped after the core result slides.
''' 
import os, shutil, csv
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from PIL import Image

# ---------------------------------------------------------------
# Paths (adjust if moving the repo)
MASTER = r'D:/Bioinformatics/datasets/reference/Deepanshu_Master_slide.pptx'
OUT = r'D:/Bioinformatics/projects/NeuralTF/docs/NeuralTF_SOTA_Final.pptx'
FIG_ROOT = r'D:/Bioinformatics/projects/NeuralTF/figures'
TOP10_CSV = r'D:/Bioinformatics/projects/NeuralTF/results/top10_neural_tfs_prioritized.csv'

# ---------------------------------------------------------------
# Theme colours – same as master (hex values converted to RGB)
BLUE = RGBColor(0x44, 0x72, 0xC4)      # 4472C4
ORANGE = RGBColor(0xED, 0x7D, 0x31)    # ED7D31
GREEN = RGBColor(0x70, 0xAD, 0x47)     # 70AD47
LIGHTBLUE = RGBColor(0x5B, 0x9B, 0xD5) # 5B9BD5
DARKSLATE = RGBColor(0x44, 0x54, 0x6A) # 44546A
GRAY_DARK = RGBColor(0x33, 0x33, 0x33) # 333333
GRAY_LIGHT = RGBColor(0xE7, 0xE6, 0xE6) # E7E6E6

# ---------------------------------------------------------------
# Font sizes (pts)
SIZE_TITLE = Pt(23)
SIZE_BODY = Pt(20)
SIZE_FOOTER = Pt(12)
SIZE_SMALL = Pt(14)

# ---------------------------------------------------------------
def add_chrome(slide, title_text=None):
    '''Insert master‑style chrome (title bar, two thin dividers, slide‑number placeholder, and empty footer).'''
    if title_text:
        tb = slide.shapes.add_textbox(Emu(0.35*914400), Emu(0.26*914400), Emu(11.83*914400), Emu(0.50*914400))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Calibri'; p.font.size = SIZE_TITLE; p.font.bold = True; p.font.color.rgb = DARKSLATE
    # horizontal dividers – top & bottom of content area
    for y in (0.76, 7.03):
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                         Emu(-0.03*914400), Emu(y*914400),
                                         Emu(13.40*914400), Emu(y*914400))
        line.line.width = Pt(1); line.line.color.rgb = GRAY_DARK
    # slide‑number placeholder (right‑aligned, bottom‑right)
    sn = slide.shapes.add_textbox(Emu(9.42*914400), Emu(6.95*914400), Emu(3.00*914400), Emu(0.40*914400))
    p = sn.text_frame.paragraphs[0]
    p.font.name = 'Calibri'; p.font.size = Pt(12); p.font.color.rgb = GRAY_DARK; p.alignment = PP_ALIGN.RIGHT; p.text = ''
    # footer – left side (empty for now)
    foot = slide.shapes.add_textbox(Emu(0.48*914400), Emu(7.04*914400), Emu(13.11*914400), Emu(0.30*914400))
    ft = foot.text_frame.paragraphs[0]
    ft.text = ''
    ft.font.name = 'Calibri'; ft.font.size = SIZE_FOOTER; ft.font.color.rgb = GRAY_DARK
    return sn

def add_paragraph(slide, left, top, w, h, lines, size=SIZE_BODY, bold=False):
    tb = slide.shapes.add_textbox(Emu(left*914400), Emu(top*914400), Emu(w*914400), Emu(h*914400))
    tf = tb.text_frame; tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.name = 'Calibri'; p.font.size = size; p.font.bold = bold

def add_picture_fit(slide, img_path, left, top, max_w, max_h):
    if not os.path.exists(img_path):
        return
    with Image.open(img_path) as im:
        w_px, h_px = im.size
    max_w_emu = Emu(max_w*914400)
    max_h_emu = Emu(max_h*914400)
    scale = min(max_w_emu / w_px, max_h_emu / h_px)
    w_emu = int(w_px * scale)
    h_emu = int(h_px * scale)
    x_emu = Emu(left*914400) + (max_w_emu - w_emu)//2
    y_emu = Emu(top*914400) + (max_h_emu - h_emu)//2
    slide.shapes.add_picture(img_path, x_emu, y_emu, w_emu, h_emu)

# ---------------------------------------------------------------
# Copy master template and start a fresh presentation
shutil.copy2(MASTER, OUT)
prs = Presentation(OUT)
# Remove any existing slides (preserve master layouts)
for sid in list(prs.slides._sldIdLst):
    rid = sid.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rid)
    prs.slides._sldIdLst.remove(sid)
blank = prs.slide_masters[0].slide_layouts[6]
title_layout = prs.slide_masters[0].slide_layouts[0]
slide_numbers = []  # collection of slide‑number placeholders

# ---------------------------------------------------------------
# 1) Title slide
s = prs.slides.add_slide(title_layout)
s.shapes.title.text = 'NeuralTF – State-of-the-Art Overview'
s.placeholders[1].text = ('Deepanshu | 5th Year BS-MS (MS-Thesis) | IISER Kolkata | '
                         'Email: d22ms218[at]iiserkol[dot]ac[dot]in')
slide_numbers.append(add_chrome(s))

# ---------------------------------------------------------------
# 2) Outline slide – short agenda
s = prs.slides.add_slide(blank)
add_paragraph(s,0.35,0.90,12,0.5,['Outline:'], size=SIZE_TITLE, bold=True)
outline = [
    '1️⃣ Introduction & biology context',
    '2️⃣ Data landscape (3 atlases + RNAi)',
    '3️⃣ Computational pipeline (3 phases, 7 evidence streams)',
    '4️⃣ Core results – score distributions, candidate funnel, top-10 dual-track',
    '5️⃣ Evidence visualisations (heatmaps, composition, stream ablation)',
    '6️⃣ Dirichlet robustness & weight-sensitivity analysis',
    '7️⃣ Conclusions & next steps'
]
add_paragraph(s,0.35,1.30,12,5,outline,size=SIZE_BODY)
slide_numbers.append(add_chrome(s, title_text='Outline'))

# ---------------------------------------------------------------
# Helper to create a figure slide (title, image, optional caption)
def add_figure_slide(fig_file, slide_title, caption_lines=None):
    s = prs.slides.add_slide(blank)
    add_paragraph(s,0.35,0.90,12,0.5,[slide_title], size=SIZE_TITLE, bold=True)
    # place image centred, max width 12in, max height 5.5in (content area)
    add_picture_fit(s, os.path.join(FIG_ROOT, fig_file), left=0.35, top=1.6, max_w=12, max_h=5.5)
    if caption_lines:
        # caption placed below image (starting at y≈7.0in, small font)
        add_paragraph(s,0.35,7.0,12,0.6,caption_lines, size=SIZE_SMALL)
    slide_numbers.append(add_chrome(s, title_text=slide_title))

# ---------------------------------------------------------------
# 3) Core result figures (one per slide)
add_figure_slide('1_score_distributions.png', 'Score distributions',
    ['Distribution of TF scores across three atlases (Fincher, Plass, King).'])
add_figure_slide('2_candidate_summary.png', 'Candidate summary',
    ['Summary of candidate TFs after initial filtering (expression & neural enrichment).'])
add_figure_slide('3_top10_dual_track.png', 'Dual-track top-10 candidates',
    ['Left-hand side: RNAi-validated TFs; Right-hand side: novel high-scoring TFs.'])
add_figure_slide('4_evidence_heatmap.png', 'Evidence heatmap (all 249 candidates)',
    ['Rows = TFs, columns = 7 evidence streams; colour encodes weighted score.'])
add_figure_slide('5_candidate_funnel.png', 'Candidate funnel',
    ['Filters: 2,800 TF catalogue → 249 expression-positive → 96 neural-enriched → 99 final set.'])
add_figure_slide('6_evidence_composition.png', 'Evidence composition',
    ['Proportion of total evidence contributed by each stream.'])
add_figure_slide('7_stream_ablation.png', 'Stream ablation analysis',
    ['Impact on candidate ranking when each evidence stream is removed in turn.'])
add_figure_slide('8_top10_radar.png', 'Top-10 radar plot',
    ['Radar visualisation of the seven evidence dimensions for the top-10 TFs.'])
add_figure_slide('9_go_dotplot.png', 'GO enrichment dot plot',
    ['Gene-ontology terms enriched among the top-10 neural TF candidates.'])

# ---------------------------------------------------------------
# 4) Dirichlet robustness & extra analyses (selected key figures)
add_figure_slide('fig_dirichlet_score_density.png', 'Dirichlet score density',
    ['Score distribution under Dirichlet-centered prior (k=40).'])
add_figure_slide('fig_dirichlet_scatter.png', 'Dirichlet scatter',
    ['Scatter of centered vs uniform scores for each candidate.'])
add_figure_slide('fig_dirichlet_trackA_top5.png', 'Dirichlet – Track A top-5',
    ['Top-5 TFs by centered Dirichlet weighting (Track A).'])
add_figure_slide('fig_dirichlet_trackB_top5.png', 'Dirichlet – Track B top-5',
    ['Top-5 TFs by uniform Dirichlet weighting (Track B).'])
add_figure_slide('fig_dirichlet_uniform_vs_centered.png', 'Uniform vs Centered comparison',
    ['Overlap of top-10 TFs between uniform and centered Dirichlet priors.'])

# ---------------------------------------------------------------
# 5) Final slide – conclusions
s = prs.slides.add_slide(blank)
add_paragraph(s,0.35,0.90,12,0.5,['Conclusions & next steps'], size=SIZE_TITLE, bold=True)
conclusion = [
    '• A robust, multi-evidence pipeline identifies high-confidence neural TFs.',
    '• Dual-track top-10 includes 5 RNAi-validated and 5 novel candidates for experimental follow-up.',
    '• Dirichlet priors demonstrate stability of rankings under different uncertainty models.',
    '• Next: wet-lab validation (dsRNA knock-down, FISH) of novel TFs; integrate additional atlases when available.'
]
add_paragraph(s,0.35,1.30,12,5,conclusion,size=SIZE_BODY)
slide_numbers.append(add_chrome(s, title_text='Conclusions'))

# ---------------------------------------------------------------
# Populate slide numbers (right-aligned bottom-right)
for idx, sn_box in enumerate(slide_numbers, start=1):
    para = sn_box.text_frame.paragraphs[0]
    para.text = str(idx)
    para.font.name = 'Calibri'; para.font.size = Pt(12); para.font.color.rgb = GRAY_DARK; para.alignment = PP_ALIGN.RIGHT

# Save the presentation
prs.save(OUT)
